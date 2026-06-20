
from lxml import etree
import zipfile, os, tempfile, shutil, re, json, csv
from rich.progress import track, Progress
from conf import *
from utils import *
from typing import Callable
from pathlib import Path

# Optional helpers
try:
	from docx import Document
except ImportError:
	Document = None

try:
	from openpyxl import load_workbook
except ImportError:
	load_workbook = None

try:
	from pptx import Presentation
except ImportError:
	Presentation = None


def translate_pdf(path: Path, translate_fn: Callable[[str], str], verbose: bool=False):
	from pypdf import PdfReader
	from reportlab.lib.pagesizes import A4
	from reportlab.pdfgen import canvas
	from reportlab.lib.units import mm
	from reportlab.pdfbase.pdfmetrics import stringWidth


	reader = PdfReader(str(path))

	tmp_out = str(path) + ".translated.pdf"
	c = canvas.Canvas(tmp_out, pagesize=A4)
	width, height = A4

	margin = 15 * mm
	max_width = width - 2 * margin
	y_start = height - margin
	font_name = "Helvetica"
	font_size = 10
	line_height = 12

	def wrap_pdf_text(text):
		lines = []
		for para in text.splitlines():
			if not para.strip():
				lines.append("")
				continue
			words = para.split()
			current = ""
			for w in words:
				test = (current + " " + w).strip()
				if stringWidth(test, font_name, font_size) <= max_width:
					current = test
				else:
					if current:
						lines.append(current)
					current = w
			if current:
				lines.append(current)
		return lines

	for i, page in enumerate(track(reader.pages, description="Translating PDF pages")):
		raw_text = page.extract_text() or ""
		raw_text = raw_text.strip()

		if verbose and raw_text:
			print(f"\n--- Page {i + 1} original ---\n{raw_text}\n")

		translated = translate_fn(raw_text) if raw_text else ""

		if verbose and translated.strip():
			print(f"\n--- Page {i + 1} translated ---\n{translated}\n")

		c.setFont(font_name, font_size)
		y = y_start

		if translated.strip():
			for line in wrap_pdf_text(translated):
				if y < margin:
					c.showPage()
					c.setFont(font_name, font_size)
					y = y_start
				c.drawString(margin, y, line)
				y -= line_height

		c.showPage()

	c.save()
	shutil.move(tmp_out, path)

# ---------------------------
# ODT / ODS / ODP
# ---------------------------
def edit_ooo_zip_inplace(path, translate_fn, tags, content_name="content.xml", verbose: bool=False):
    with zipfile.ZipFile(path, 'r') as z:
        names = z.namelist()
        if content_name not in names:
            raise ValueError(f"{content_name} not found in archive")
        content = z.read(content_name)
        other_files = {name: z.read(name) for name in names if name != content_name}

    root = etree.fromstring(content)

    paras = []
    for tag in tags:
        paras.extend(root.findall(f'.//{tag}', NS))

    for p in track(paras):
        orig = paragraph_text(p).strip()
        if orig and clean_text(REG_CLEAN, orig).strip():
            set_mixed_text(p, translate_text(orig, translate_fn, verbose=verbose))

    new_content = etree.tostring(root, xml_declaration=True, encoding='UTF-8', standalone=False)

    tmp_fd, tmp_path = tempfile.mkstemp(suffix='.zip')
    os.close(tmp_fd)
    try:
        with zipfile.ZipFile(tmp_path, 'w', compression=zipfile.ZIP_DEFLATED) as z:
            z.writestr(content_name, new_content)
            for name, data in other_files.items():
                z.writestr(name, data)
        shutil.move(tmp_path, path)
    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)


def translate_odt(path: Path, translate_fn: Callable[[str], str], verbose: bool=False):
    edit_ooo_zip_inplace(path, translate_fn, tags=('text:p', 'text:h'), verbose=verbose)


def translate_ods(path: Path, translate_fn: Callable[[str], str], verbose: bool=False):
    edit_ooo_zip_inplace(path, translate_fn, tags=('text:p',), content_name="content.xml", verbose=verbose)


def translate_odp(path: Path, translate_fn: Callable[[str], str], verbose: bool=False):
    edit_ooo_zip_inplace(path, translate_fn, tags=('text:p', 'text:h'), verbose=verbose)


# ---------------------------
# DOCX
# ---------------------------
def translate_docx(path: Path, translate_fn: Callable[[str], str], verbose: bool=False):
	if Document is None:
		raise ImportError("python-docx is required for .docx files")
	doc = Document(str(path))

	if doc.paragraphs:
		for p in track(doc.paragraphs, "Paragraphs..."):
			txt = p.text.strip()
			if txt and clean_text(REG_CLEAN, txt).strip():
				p.text = translate_text(txt, translate_fn, verbose=verbose)
	
	if doc.tables:
		tables_len = len(doc.tables)
		for i in range(tables_len):
			table = doc.tables[i]
			for row in track(table.rows, description=f"Table {i}/{tables_len}"):
				for cell in row.cells:
					for p in cell.paragraphs:
						txt = p.text.strip()
						if txt and clean_text(REG_CLEAN, txt).strip():
							p.text = translate_text(txt, translate_fn, verbose=verbose)

	doc.save(str(path))


# ---------------------------
# XLSX
# ---------------------------
def translate_xlsx(path: Path, translate_fn: Callable[[str], str], verbose: bool=False):
	if load_workbook is None:
		raise ImportError("openpyxl is required for .xlsx files")
	wb = load_workbook(str(path))

	worksheets_len = len(wb.worksheets)
	for i in range(worksheets_len):
		rows = wb.worksheets[i].iter_rows()
		for row in track(rows, description=f"Table {i}/{worksheets_len}"):
			for cell in row:
				if isinstance(cell.value, str) and cell.value.strip():
					if clean_text(REG_CLEAN, cell.value).strip():
						cell.value = translate_text(cell.value, translate_fn, verbose=verbose)

	wb.save(str(path))


# ---------------------------
# PPTX
# ---------------------------
def translate_pptx(path: Path, translate_fn: Callable[[str], str], verbose: bool=False):
	if Presentation is None:
		raise ImportError("python-pptx is required for .pptx files")
	prs = Presentation(str(path))

	def translate_shape(shape):
		if hasattr(shape, "text_frame") and shape.has_text_frame:
			for p in shape.text_frame.paragraphs:
				for run in p.runs:
					txt = run.text.strip()
					if txt and clean_text(REG_CLEAN, txt).strip():
						run.text = translate_text(txt, translate_fn, verbose=verbose)

	for slide in track(prs.slides):
		for shape in slide.shapes:
			translate_shape(shape)

	prs.save(str(path))


